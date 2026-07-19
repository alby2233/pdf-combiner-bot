import os
import sys
import tempfile
import io

def word_to_pdf(doc_path, pdf_path):
    """Convert Word document (.doc, .docx) to PDF using MS Word COM interface."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(doc_path))
        # 17 represents wdFormatPDF
        doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)
        doc.Close()
    finally:
        try:
            word.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

def pdf_to_word(pdf_path, doc_path):
    """Convert PDF to Word document (.docx) using pdf2docx library."""
    from pdf2docx import Converter
    cv = Converter(pdf_path)
    try:
        cv.convert(doc_path, start=0, end=None)
    finally:
        cv.close()



def ppt_to_pdf(ppt_path, pdf_path):
    """Convert PowerPoint presentation (.ppt, .pptx) to PDF using MS PowerPoint COM interface."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    try:
        ppt = win32com.client.DispatchEx("PowerPoint.Application")
        # WithWindow=False runs ppt in background
        pres = ppt.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)
        # 32 represents ppSaveAsPDF
        pres.SaveAs(os.path.abspath(pdf_path), FileFormat=32)
        pres.Close()
    finally:
        try:
            ppt.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

def pdf_to_ppt(pdf_path, ppt_path):
    """Convert PDF to PowerPoint presentation (.pptx) by converting pages to images and adding them to slides."""
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Pt
    
    prs = Presentation()
    # Remove default slides
    for i in range(len(prs.slides)-1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    doc = fitz.open(pdf_path)
    try:
        # Create a temp folder for slide images
        with tempfile.TemporaryDirectory() as temp_dir:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                img_temp_path = os.path.join(temp_dir, f"slide_page_{page_num}.png")
                pix.save(img_temp_path)
                
                rect = page.rect
                width_pt = rect.width
                height_pt = rect.height
                
                # Set slide size to match page dimensions (for the first page)
                if page_num == 0:
                    prs.slide_width = Pt(width_pt)
                    prs.slide_height = Pt(height_pt)
                
                blank_slide_layout = prs.slide_layouts[6] # index 6 is blank slide layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add image covering the entire slide
                slide.shapes.add_picture(img_temp_path, Pt(0), Pt(0), width=Pt(width_pt), height=Pt(height_pt))
                
            prs.save(ppt_path)
    finally:
        doc.close()

def excel_to_pdf(xls_path, pdf_path):
    """Convert Excel spreadsheet (.xls, .xlsx) to PDF using MS Excel COM interface."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        wb = excel.Workbooks.Open(os.path.abspath(xls_path))
        # 0 represents xlTypePDF
        wb.ExportAsFixedFormat(0, os.path.abspath(pdf_path))
        wb.Close(SaveChanges=False)
    finally:
        try:
            excel.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

def pdf_to_excel(pdf_path, xls_path):
    """Convert PDF to Excel workbook (.xlsx) by extracting tables or text using pdfplumber."""
    import pdfplumber
    import openpyxl
    from openpyxl import Workbook
    
    wb = Workbook()
    # Remove default active sheet
    default_sheet = wb.active
    
    with pdfplumber.open(pdf_path) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page {page_idx + 1}")
            tables = page.extract_tables()
            
            if tables:
                row_cursor = 1
                for table in tables:
                    for row in table:
                        for col_idx, cell in enumerate(row):
                            ws.cell(row=row_cursor, column=col_idx + 1, value=cell)
                        row_cursor += 1
                    row_cursor += 2  # spacing between tables on the same sheet
            else:
                text = page.extract_text()
                if text:
                    for row_idx, line in enumerate(text.split('\n')):
                        ws.cell(row=row_idx + 1, column=1, value=line)
                        
    if len(wb.sheetnames) > 1 and "Sheet" in wb.sheetnames:
        wb.remove(wb["Sheet"])
    elif len(wb.sheetnames) == 0:
        wb.create_sheet(title="Empty")
        
    wb.save(xls_path)

def ppt_to_images(ppt_path, output_dir):
    """Export PowerPoint slides as images to a directory using MS PowerPoint COM interface."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    try:
        ppt = win32com.client.DispatchEx("PowerPoint.Application")
        pres = ppt.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)
        
        os.makedirs(output_dir, exist_ok=True)
        image_paths = []
        
        for idx, slide in enumerate(pres.Slides):
            img_path = os.path.join(output_dir, f"slide_{idx+1}.png")
            # Export as PNG
            slide.Export(os.path.abspath(img_path), "PNG")
            image_paths.append(img_path)
            
        pres.Close()
        return image_paths
    finally:
        try:
            ppt.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

def pdf_to_images(pdf_path, output_dir):
    """Convert PDF pages to PNG images using PyMuPDF."""
    import fitz
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_paths = []
    try:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(output_dir, f"page_{page_num+1}.png")
            pix.save(img_path)
            image_paths.append(img_path)
        return image_paths
    finally:
        doc.close()

def images_to_pdf(image_paths, pdf_path):
    """Convert a list of images into a single PDF file using Pillow."""
    from PIL import Image
    images = []
    for img_path in image_paths:
        img = Image.open(img_path)
        if img.mode != "RGB":
            img = img.convert("RGB")
        images.append(img)
    if images:
        images[0].save(pdf_path, save_all=True, append_images=images[1:])

def merge_pdfs(file_paths, output_path):
    """Merge multiple PDF files into a single PDF using pypdf."""
    from pypdf import PdfWriter
    writer = PdfWriter()
    for path in file_paths:
        writer.append(path)
    with open(output_path, "wb") as f:
        writer.write(f)
    writer.close()

def split_pdf(pdf_path, page_ranges_str, output_path):
    """Extract page ranges (e.g. '1, 3-5, 8') from a PDF using pypdf."""
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    
    pages_to_extract = set()
    parts = page_ranges_str.split(",")
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_str, end_str = part.split("-")
            try:
                start = int(start_str.strip())
                end = int(end_str.strip())
                for p in range(start, end + 1):
                    if 1 <= p <= total_pages:
                        pages_to_extract.add(p - 1)
            except ValueError:
                continue
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages_to_extract.add(p - 1)
            except ValueError:
                continue
                
    for p in sorted(pages_to_extract):
        writer.add_page(reader.pages[p])
        
    with open(output_path, "wb") as f:
        writer.write(f)

def rotate_pdf(pdf_path, angle, output_path):
    """Rotate all pages of a PDF by a given angle (90, 180, 270) using pypdf."""
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def add_header_footer_page_numbers(pdf_path, output_path, header_text=None, footer_text=None, add_page_numbers=True, start_page_num=1, exclude_first_page=False):
    """Apply customized header, footer, and page numbers onto a PDF using reportlab overlay."""
    import io
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    
    for i in range(total_pages):
        page = reader.pages[i]
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        
        has_header = bool(header_text)
        has_footer = bool(footer_text)
        has_page_num = bool(add_page_numbers)
        is_excluded = (i == 0 and exclude_first_page)
        
        if (has_header or has_footer or has_page_num) and not is_excluded:
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=(width, height))
            can.setFont("Helvetica", 9)
            can.setFillColorRGB(0.3, 0.3, 0.3)
            
            margin = 54 # 0.75 inch
            
            if has_header:
                can.drawString(margin, height - 36, header_text)
                can.setStrokeColorRGB(0.8, 0.8, 0.8)
                can.setLineWidth(0.5)
                can.line(margin, height - 42, width - margin, height - 42)
                
            if has_footer:
                can.drawString(margin, 36, footer_text)
                
            if has_page_num:
                page_str = f"{i + start_page_num}"
                can.drawRightString(width - margin, 36, page_str)
                
            can.save()
            packet.seek(0)
            overlay_reader = PdfReader(packet)
            page.merge_page(overlay_reader.pages[0])
            
        writer.add_page(page)
        
    with open(output_path, "wb") as f:
        writer.write(f)
