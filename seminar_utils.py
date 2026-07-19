import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

THEMES = {
    "dark": {
        "bg_color": RGBColor(20, 24, 33),
        "title_color": RGBColor(0, 220, 255),
        "text_color": RGBColor(240, 240, 240),
        "font_name": "Trebuchet MS",
        # Docx equivalents
        "docx_bg": DocxRGBColor(20, 24, 33),
        "docx_title": DocxRGBColor(0, 160, 200),
        "docx_text": DocxRGBColor(50, 50, 50), # Word documents are typically light mode for printing, but we style titles/headings
        "docx_font": "Segoe UI"
    },
    "clean": {
        "bg_color": RGBColor(255, 255, 255),
        "title_color": RGBColor(30, 30, 30),
        "text_color": RGBColor(80, 80, 80),
        "font_name": "Arial",
        "docx_bg": DocxRGBColor(255, 255, 255),
        "docx_title": DocxRGBColor(30, 30, 30),
        "docx_text": DocxRGBColor(80, 80, 80),
        "docx_font": "Arial"
    },
    "corporate": {
        "bg_color": RGBColor(245, 247, 250),
        "title_color": RGBColor(10, 35, 80),
        "text_color": RGBColor(50, 60, 70),
        "font_name": "Calibri",
        "docx_bg": DocxRGBColor(245, 247, 250),
        "docx_title": DocxRGBColor(10, 35, 80),
        "docx_text": DocxRGBColor(50, 60, 70),
        "docx_font": "Calibri"
    },
    "warm": {
        "bg_color": RGBColor(250, 244, 230),
        "title_color": RGBColor(110, 45, 10),
        "text_color": RGBColor(70, 50, 40),
        "font_name": "Georgia",
        "docx_bg": DocxRGBColor(250, 244, 230),
        "docx_title": DocxRGBColor(110, 45, 10),
        "docx_text": DocxRGBColor(70, 50, 40),
        "docx_font": "Georgia"
    }
}

def create_seminar_presentation(output_path, title, slides_data, theme_key="clean"):
    theme = THEMES.get(theme_key, THEMES["clean"])
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    # 1. Generate Title Slide
    slide = prs.slides.add_slide(blank_layout)
    
    # Set background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["bg_color"]
    bg.line.fill.background()
    
    # Large Title text box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = theme["font_name"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    
    p2 = tf.add_paragraph()
    p2.text = "Seminar Presentation & Reference Notes"
    p2.font.name = theme["font_name"]
    p2.font.size = Pt(22)
    p2.font.color.rgb = theme["text_color"]
    p2.space_before = Pt(15)
    
    # 2. Generate Content Slides
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme["bg_color"]
        bg.line.fill.background()
        
        # Header Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.333), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = theme["title_color"]
        line.line.fill.background()
        
        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Topic Overview")
        p.font.name = theme["font_name"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = theme["title_color"]
        
        # Slide Bullets
        bullets = slide_data.get("bullets", [])
        b_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        tf_bullets = b_box.text_frame
        tf_bullets.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p_bullet = tf_bullets.add_paragraph() if i > 0 else tf_bullets.paragraphs[0]
            p_bullet.text = bullet
            p_bullet.font.name = theme["font_name"]
            p_bullet.font.size = Pt(20)
            p_bullet.font.color.rgb = theme["text_color"]
            p_bullet.space_after = Pt(14)
            
        # Add Speaker Notes
        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            tf_notes = notes_slide.notes_text_frame
            tf_notes.text = notes
            
    prs.save(output_path)

def create_seminar_report(output_path, title, report_data, theme_key="clean"):
    theme = THEMES.get(theme_key, THEMES["clean"])
    
    doc = Document()
    
    # Styles configuration
    style_normal = doc.styles['Normal']
    style_normal.font.name = theme["docx_font"]
    style_normal.font.size = DocxPt(11)
    style_normal.font.color.rgb = theme["docx_text"]
    
    # Cover Page
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_before = DocxPt(120)
    run_title = title_p.add_run(title.upper())
    run_title.font.name = theme["docx_font"]
    run_title.font.size = DocxPt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = theme["docx_title"]
    
    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_p.paragraph_format.space_after = DocxPt(180)
    run_sub = sub_p.add_run("SEMINAR PRESENTATION & RESEARCH REPORT")
    run_sub.font.name = theme["docx_font"]
    run_sub.font.size = DocxPt(14)
    run_sub.font.color.rgb = DocxRGBColor(120, 120, 120)
    
    info_p = doc.add_paragraph()
    info_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_info = info_p.add_run("Prepared for Seminar Presentation\nGenerated via Academic AI Engine")
    run_info.font.italic = True
    run_info.font.size = DocxPt(10)
    
    doc.add_page_break()
    
    # Report Content Sections
    for section in report_data:
        heading_text = section.get("heading", "")
        content_text = section.get("content", "")
        
        if not heading_text:
            continue
            
        h = doc.add_heading(level=1)
        h.paragraph_format.space_before = DocxPt(18)
        h.paragraph_format.space_after = DocxPt(6)
        
        h_run = h.add_run(heading_text)
        h_run.font.name = theme["docx_font"]
        h_run.font.size = DocxPt(16)
        h_run.font.bold = True
        h_run.font.color.rgb = theme["docx_title"]
        
        # Split content by newlines to form clean paragraphs
        paragraphs = content_text.split("\n")
        for para in paragraphs:
            para = para.strip()
            if para:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = DocxPt(8)
                p.paragraph_format.line_spacing = 1.15
                p_run = p.add_run(para)
                p_run.font.name = theme["docx_font"]
                p_run.font.size = DocxPt(11)
                p_run.font.color.rgb = theme["docx_text"]
                
    doc.save(output_path)
